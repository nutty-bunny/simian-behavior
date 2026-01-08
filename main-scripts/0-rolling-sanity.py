import os
import pickle
from pathlib import Path
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import matplotlib.dates as mdates
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import io

# ============================================================================
# CONFIGURATION
# ============================================================================
directory = "/Users/similovesyou/Desktop/qts/simian-behavior/data/py"
base = "/Users/similovesyou/Desktop/qts/simian-behavior/5-csrt"
pickle_rick = os.path.join(directory, "data.pickle")

plt.rcParams["svg.fonttype"] = "none"

# ============================================================================
# LOAD DATA
# ============================================================================
with open(pickle_rick, "rb") as handle:
    data = pickle.load(handle)

# Remove monkeys with RFID issues
monkeys_to_remove = {"joy", "jipsy"}
for species in data:
    for name in list(data[species].keys()):
        if name in monkeys_to_remove:
            del data[species][name]

# ============================================================================
# HIERARCHY DATA
# ============================================================================
tonkean_hierarchy_file = os.path.join(directory, "hierarchy/tonkean/elo_matrix.xlsx")
rhesus_hierarchy_file = os.path.join(directory, "hierarchy/rhesus/elo_matrix.xlsx")

tonkean_abrv_2_fn = {
    "abr": "abricot", "ala": "alaryc", "alv": "alvin", "anu": "anubis",
    "bar": "barnabe", "ber": "berenice", "ces": "cesar", "dor": "dory",
    "eri": "eric", "jea": "jeanne", "lad": "lady", "las": "lassa",
    "nem": "nema", "nen": "nenno", "ner": "nereis", "ola": "olaf",
    "olg": "olga", "oli": "olli", "pac": "patchouli", "pat": "patsy",
    "wal": "wallace", "wat": "walt", "wot": "wotan", "yan": "yang",
    "yak": "yannick", "yin": "yin", "yoh": "yoh", "ult": "ulysse",
    "fic": "ficelle", "gan": "gandhi", "hav": "havanna", "con": "controle",
    "gai": "gaia", "her": "hercules", "han": "hanouk", "hor": "horus",
    "imo": "imoen", "ind": "indigo", "iro": "iron", "isi": "isis",
    "joy": "joy", "jip": "jipsy",
}

rhesus_abrv_2_fn = {
    "the": "theoden", "any": "anyanka", "djo": "djocko", "vla": "vladimir",
    "yel": "yelena", "bor": "boromir", "far": "faramir", "yva": "yvan",
    "baa": "baal", "spl": "spliff", "ol": "ol", "nat": "natasha",
    "arw": "arwen", "eow": "eowyn", "mar": "mar"
}

tonkean_hierarchy = pd.read_excel(tonkean_hierarchy_file)
rhesus_hierarchy = pd.read_excel(rhesus_hierarchy_file)

tonkean_hierarchy = tonkean_hierarchy.rename(columns=tonkean_abrv_2_fn)
rhesus_hierarchy = rhesus_hierarchy.rename(columns=rhesus_abrv_2_fn)

hierarchies = {
    "tonkean": tonkean_hierarchy,
    "rhesus": rhesus_hierarchy
}

# ============================================================================
# BUILD ELO DICTIONARY
# ============================================================================
elo = {}
for species, hierarchy in hierarchies.items():
    hierarchy["Date"] = pd.to_datetime(hierarchy["Date"])
    for name in hierarchy.columns[1:]:
        if name in hierarchy:
            monkey_elo = hierarchy[["Date", name]].dropna()
            if not monkey_elo.empty:
                first_date = monkey_elo["Date"].min()
                cutoff_date = first_date + pd.DateOffset(months=2)
                filtered_elo = monkey_elo[monkey_elo["Date"] > cutoff_date]
                elo[name] = filtered_elo

# ============================================================================
# COMPUTE SESSION PROPORTIONS
# ============================================================================
def compute_session_proportions(species, name, monkey_data, elo_dict):
    """Compute session-based performance proportions."""
    attempts = monkey_data["attempts"].copy()
    
    attempts["instant_begin"] = pd.to_datetime(attempts["instant_begin"], unit="ms", errors="coerce")
    attempts["instant_end"] = pd.to_datetime(attempts["instant_end"], unit="ms", errors="coerce")
    
    # Filter by ELO period if available
    if name in elo_dict and not elo_dict[name].empty:
        first_elo = elo_dict[name]["Date"].min()
        last_elo = elo_dict[name]["Date"].max()
        
        attempts = attempts[
            (attempts["instant_begin"] >= first_elo) & 
            (attempts["instant_end"] <= last_elo)
        ]
    
    if attempts.empty:
        return pd.DataFrame()
    
    # Group by session
    session_data = []
    for session, session_attempts in attempts.groupby("session"):
        total_attempts = len(session_attempts)
        outcome_counts = session_attempts["result"].value_counts().to_dict()
        
        session_data.append({
            "session": session,
            "p_success": outcome_counts.get("success", 0.0) / total_attempts,
            "p_error": outcome_counts.get("error", 0.0) / total_attempts,
            "p_omission": outcome_counts.get("stepomission", 0.0) / total_attempts,
            "p_premature": outcome_counts.get("prematured", 0.0) / total_attempts,
            "total_attempts": total_attempts,
        })
    
    session_df = pd.DataFrame(session_data)
    session_df = session_df.sort_values("session").reset_index(drop=True)
    
    return session_df

# ============================================================================
# COMPUTE DAILY PROPORTIONS
# ============================================================================
def compute_daily_proportions(species, name, monkey_data, elo_dict):
    """Compute daily performance proportions - only days with data."""
    attempts = monkey_data["attempts"].copy()
    
    attempts["instant_begin"] = pd.to_datetime(attempts["instant_begin"], unit="ms", errors="coerce")
    attempts["instant_end"] = pd.to_datetime(attempts["instant_end"], unit="ms", errors="coerce")
    attempts["date"] = attempts["instant_begin"].dt.date
    
    # Filter by ELO period if available
    if name in elo_dict and not elo_dict[name].empty:
        first_elo = elo_dict[name]["Date"].min().date()
        last_elo = elo_dict[name]["Date"].max().date()
        
        attempts = attempts[
            (attempts["date"] >= first_elo) & 
            (attempts["date"] <= last_elo)
        ]
    
    if attempts.empty:
        return pd.DataFrame()
    
    # Group by date
    daily_data = []
    day_counter = 1
    
    for date, day_attempts in attempts.groupby("date"):
        total_attempts = len(day_attempts)
        outcome_counts = day_attempts["result"].value_counts().to_dict()
        
        daily_data.append({
            "date": pd.to_datetime(date),
            "day_number": day_counter,
            "p_success": outcome_counts.get("success", 0.0) / total_attempts,
            "p_error": outcome_counts.get("error", 0.0) / total_attempts,
            "p_omission": outcome_counts.get("stepomission", 0.0) / total_attempts,
            "p_premature": outcome_counts.get("prematured", 0.0) / total_attempts,
            "total_attempts": total_attempts,
        })
        
        day_counter += 1
    
    daily_df = pd.DataFrame(daily_data)
    daily_df = daily_df.sort_values("date").reset_index(drop=True)
    
    return daily_df

# ============================================================================
# COLLECT DATA FOR ALL MONKEYS
# ============================================================================
all_monkeys = {}

print("\nCollecting data for all monkeys...")
for species in ["tonkean", "rhesus"]:
    if species not in data:
        continue
        
    for name, monkey_data in data[species].items():
        session_df = compute_session_proportions(species, name, monkey_data, elo)
        daily_df = compute_daily_proportions(species, name, monkey_data, elo)
        
        if not session_df.empty and not daily_df.empty:
            all_monkeys[name] = {
                "species": species,
                "session_df": session_df,
                "daily_df": daily_df
            }
            print(f"  {name}: {len(session_df)} sessions, {len(daily_df)} days with data")

print(f"\nTotal monkeys: {len(all_monkeys)}")

# ============================================================================
# HELPER FUNCTION TO PLOT WITH GAPS
# ============================================================================
def plot_stacked_with_gaps(ax, dates, data_dict, colors, labels, max_gap_days=1):
    """
    Plot stacked area with breaks when there are gaps larger than max_gap_days.
    
    Parameters:
    - ax: matplotlib axis
    - dates: array of datetime objects
    - data_dict: dict with keys matching labels, values are arrays
    - colors: list of colors for each data series
    - labels: list of labels for each data series
    - max_gap_days: maximum days between consecutive points before inserting gap
    """
    
    # Find gaps larger than max_gap_days
    date_diffs = np.diff(dates)
    gap_indices = np.where(date_diffs > pd.Timedelta(days=max_gap_days))[0]
    
    # Split data at gaps
    split_points = [0] + [i + 1 for i in gap_indices] + [len(dates)]
    
    # Plot each continuous segment
    for i in range(len(split_points) - 1):
        start_idx = split_points[i]
        end_idx = split_points[i + 1]
        
        segment_dates = dates[start_idx:end_idx]
        
        # Only add label on first segment
        segment_labels = labels if i == 0 else [None] * len(labels)
        
        # Stack the data
        ax.stackplot(
            segment_dates,
            *[data_dict[label][start_idx:end_idx] for label in labels],
            colors=colors,
            labels=segment_labels,
            edgecolor="none",
            linewidth=0,
            baseline="zero",
            antialiased=False
        )

# ============================================================================
# PLOTTING FUNCTIONS
# ============================================================================
def create_raw_plot(monkey_name, species_name, session_df, daily_df, use_dates=False):
    """Create raw data plot (no rolling average)."""
    
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 6))
    
    # ========== LEFT: SESSION-BASED (RAW) - ALWAYS CONTINUOUS ==========
    stacked_df_session = session_df[["p_success", "p_error", "p_premature", "p_omission"]].copy()
    stacked_df_session = stacked_df_session.div(stacked_df_session.sum(axis=1), axis=0)
    
    ax1.stackplot(
        range(1, len(stacked_df_session) + 1),
        stacked_df_session["p_success"],
        stacked_df_session["p_error"],
        stacked_df_session["p_premature"],
        stacked_df_session["p_omission"],
        colors=["#5386b6", "#d57459", "#e8bc60", "#8d5993"],
        edgecolor="none", linewidth=0, baseline="zero",
        labels=["Success", "Error", "Premature", "Omission"],
        antialiased=False
    )
    
    overall_p_success = session_df["p_success"].mean()
    ax1.axhline(y=overall_p_success, color="black", linestyle=(0, (8, 4)), linewidth=1.0)
    
    ax1.set_title("Session-Based (Raw)", fontsize=18, fontweight="bold", fontname="Times New Roman")
    ax1.set_xlabel("Session Count", fontsize=14, fontname="Times New Roman")
    ax1.set_ylabel("Stacked Proportion", fontsize=14, fontname="Times New Roman")
    ax1.set_xlim(1, len(stacked_df_session))
    ax1.set_ylim(0, 1)
    ax1.set_yticks([0, 1])
    ax1.legend(loc="upper left", fontsize=10, prop={"family": "Times New Roman"})
    ax1.grid(True, linestyle="--", alpha=0.5)
    ax1.xaxis.set_major_locator(plt.MaxNLocator(nbins=10, integer=True))
    
    # ========== RIGHT: DAY-BASED (RAW) ==========
    stacked_df_daily = daily_df[["p_success", "p_error", "p_premature", "p_omission"]].copy()
    stacked_df_daily = stacked_df_daily.div(stacked_df_daily.sum(axis=1), axis=0)
    
    if use_dates:
        # Use actual dates - break at gaps
        dates = daily_df["date"].values
        data_dict = {
            "Success": stacked_df_daily["p_success"].values,
            "Error": stacked_df_daily["p_error"].values,
            "Premature": stacked_df_daily["p_premature"].values,
            "Omission": stacked_df_daily["p_omission"].values
        }
        
        plot_stacked_with_gaps(
            ax2, dates, data_dict,
            colors=["#5386b6", "#d57459", "#e8bc60", "#8d5993"],
            labels=["Success", "Error", "Premature", "Omission"],
            max_gap_days=1
        )
        
        ax2.set_xlabel("Date", fontsize=14, fontname="Times New Roman")
        ax2.set_xlim(dates.min(), dates.max())
        
        # Vertical date labels with more ticks
        ax2.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m-%d"))
        ax2.xaxis.set_major_locator(mdates.AutoDateLocator(maxticks=20))
        plt.setp(ax2.xaxis.get_majorticklabels(), rotation=90, ha='center', fontsize=9)
        
        # Vertical grid lines
        ax2.grid(True, axis='x', linestyle="--", alpha=0.5)
        ax2.grid(False, axis='y')
    else:
        # Use sequential day numbers - no gaps
        day_numbers = daily_df["day_number"].values.astype(int)
        ax2.stackplot(
            day_numbers,
            stacked_df_daily["p_success"],
            stacked_df_daily["p_error"],
            stacked_df_daily["p_premature"],
            stacked_df_daily["p_omission"],
            colors=["#5386b6", "#d57459", "#e8bc60", "#8d5993"],
            edgecolor="none", linewidth=0, baseline="zero",
            labels=["Success", "Error", "Premature", "Omission"],
            antialiased=False
        )
        ax2.set_xlabel("Days (with data)", fontsize=14, fontname="Times New Roman")
        ax2.set_xlim(1, len(daily_df))
        ax2.xaxis.set_major_locator(plt.MaxNLocator(nbins=10, integer=True))
        ax2.grid(True, linestyle="--", alpha=0.5)
    
    overall_p_success_daily = daily_df["p_success"].mean()
    ax2.axhline(y=overall_p_success_daily, color="black", linestyle=(0, (8, 4)), linewidth=1.0)
    
    ax2.set_title("Day-Based (Raw)", fontsize=18, fontweight="bold", fontname="Times New Roman")
    ax2.set_ylabel("Stacked Proportion", fontsize=14, fontname="Times New Roman")
    ax2.set_ylim(0, 1)
    ax2.set_yticks([0, 1])
    ax2.legend(loc="upper left", fontsize=10, prop={"family": "Times New Roman"})
    
    # Create subtitle
    if use_dates:
        subtitle = f"{monkey_name.capitalize()} ({species_name.capitalize()}) - Raw Data (Calendar Dates with Gaps)"
    else:
        subtitle = f"{monkey_name.capitalize()} ({species_name.capitalize()}) - Raw Data (Sequential Days)"
    
    fig.suptitle(subtitle, fontsize=20, fontweight="bold", fontname="Times New Roman", y=1.02)
    
    plt.tight_layout()
    return fig

def create_smoothed_plot(monkey_name, species_name, session_df, daily_df, use_dates=False):
    """Create smoothed plot (with rolling averages)."""
    
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(16, 6))
    
    # ========== LEFT: SESSION-BASED (SMOOTHED - 10 sessions) ==========
    session_df_smooth = session_df.copy()
    for col in ["p_success", "p_error", "p_premature", "p_omission"]:
        session_df_smooth[f"{col}_smoothed"] = session_df_smooth[col].rolling(
            window=10, min_periods=1, center=True
        ).mean()
    
    stacked_df_session = session_df_smooth[[
        "p_success_smoothed", "p_error_smoothed", "p_premature_smoothed", "p_omission_smoothed"
    ]].dropna()
    stacked_df_session = stacked_df_session.div(stacked_df_session.sum(axis=1), axis=0)
    
    ax1.stackplot(
        range(1, len(stacked_df_session) + 1),
        stacked_df_session["p_success_smoothed"],
        stacked_df_session["p_error_smoothed"],
        stacked_df_session["p_premature_smoothed"],
        stacked_df_session["p_omission_smoothed"],
        colors=["#5386b6", "#d57459", "#e8bc60", "#8d5993"],
        edgecolor="none", linewidth=0, baseline="zero",
        labels=["Success", "Error", "Premature", "Omission"],
        antialiased=False
    )
    
    overall_p_success = session_df["p_success"].mean()
    ax1.axhline(y=overall_p_success, color="black", linestyle=(0, (8, 4)), linewidth=1.0)
    
    ax1.set_title("Session-Based (10-session rolling avg)", fontsize=18, fontweight="bold", fontname="Times New Roman")
    ax1.set_xlabel("Session Count", fontsize=14, fontname="Times New Roman")
    ax1.set_ylabel("Stacked Proportion", fontsize=14, fontname="Times New Roman")
    ax1.set_xlim(1, len(stacked_df_session))
    ax1.set_ylim(0, 1)
    ax1.set_yticks([0, 1])
    ax1.legend(loc="upper left", fontsize=10, prop={"family": "Times New Roman"})
    ax1.grid(True, linestyle="--", alpha=0.5)
    ax1.xaxis.set_major_locator(plt.MaxNLocator(nbins=10, integer=True))
    
    # ========== RIGHT: DAY-BASED (SMOOTHED - 7 days) ==========
    daily_df_smooth = daily_df.copy()
    for col in ["p_success", "p_error", "p_premature", "p_omission"]:
        daily_df_smooth[f"{col}_smoothed"] = daily_df_smooth[col].rolling(
            window=7, min_periods=1, center=True
        ).mean()
    
    stacked_df_daily = daily_df_smooth[[
        "p_success_smoothed", "p_error_smoothed", "p_premature_smoothed", "p_omission_smoothed"
    ]].dropna()
    stacked_df_daily = stacked_df_daily.div(stacked_df_daily.sum(axis=1), axis=0)
    
    if use_dates:
        # Use actual dates - break at gaps
        dates = daily_df_smooth.loc[stacked_df_daily.index, "date"].values
        data_dict = {
            "Success": stacked_df_daily["p_success_smoothed"].values,
            "Error": stacked_df_daily["p_error_smoothed"].values,
            "Premature": stacked_df_daily["p_premature_smoothed"].values,
            "Omission": stacked_df_daily["p_omission_smoothed"].values
        }
        
        plot_stacked_with_gaps(
            ax2, dates, data_dict,
            colors=["#5386b6", "#d57459", "#e8bc60", "#8d5993"],
            labels=["Success", "Error", "Premature", "Omission"],
            max_gap_days=1
        )
        
        ax2.set_xlabel("Date", fontsize=14, fontname="Times New Roman")
        ax2.set_xlim(dates.min(), dates.max())
        
        # Vertical date labels with more ticks
        ax2.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m-%d"))
        ax2.xaxis.set_major_locator(mdates.AutoDateLocator(maxticks=20))
        plt.setp(ax2.xaxis.get_majorticklabels(), rotation=90, ha='center', fontsize=9)
        
        # Vertical grid lines
        ax2.grid(True, axis='x', linestyle="--", alpha=0.5)
        ax2.grid(False, axis='y')
    else:
        # Use sequential day numbers - no gaps
        day_numbers = daily_df_smooth.loc[stacked_df_daily.index, "day_number"].values.astype(int)
        ax2.stackplot(
            day_numbers,
            stacked_df_daily["p_success_smoothed"],
            stacked_df_daily["p_error_smoothed"],
            stacked_df_daily["p_premature_smoothed"],
            stacked_df_daily["p_omission_smoothed"],
            colors=["#5386b6", "#d57459", "#e8bc60", "#8d5993"],
            edgecolor="none", linewidth=0, baseline="zero",
            labels=["Success", "Error", "Premature", "Omission"],
            antialiased=False
        )
        ax2.set_xlabel("Days (with data)", fontsize=14, fontname="Times New Roman")
        ax2.set_xlim(day_numbers.min(), day_numbers.max())
        ax2.xaxis.set_major_locator(plt.MaxNLocator(nbins=10, integer=True))
        ax2.grid(True, linestyle="--", alpha=0.5)
    
    overall_p_success_daily = daily_df["p_success"].mean()
    ax2.axhline(y=overall_p_success_daily, color="black", linestyle=(0, (8, 4)), linewidth=1.0)
    
    ax2.set_title("Day-Based (7-day rolling avg)", fontsize=18, fontweight="bold", fontname="Times New Roman")
    ax2.set_ylabel("Stacked Proportion", fontsize=14, fontname="Times New Roman")
    ax2.set_ylim(0, 1)
    ax2.set_yticks([0, 1])
    ax2.legend(loc="upper left", fontsize=10, prop={"family": "Times New Roman"})
    
    # Create subtitle
    if use_dates:
        subtitle = f"{monkey_name.capitalize()} ({species_name.capitalize()}) - Smoothed Data (Calendar Dates with Gaps)"
    else:
        subtitle = f"{monkey_name.capitalize()} ({species_name.capitalize()}) - Smoothed Data (Sequential Days)"
    
    fig.suptitle(subtitle, fontsize=20, fontweight="bold", fontname="Times New Roman", y=1.02)
    
    plt.tight_layout()
    return fig

# ============================================================================
# CREATE POWERPOINT
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

downloads_dir = Path.home() / "Downloads"

print("\nGenerating slides...")

for monkey_name, monkey_info in all_monkeys.items():
    print(f"\n{'='*60}")
    print(f"Processing {monkey_name}...")
    print(f"{'='*60}")
    
    species_name = monkey_info["species"]
    session_df = monkey_info["session_df"].copy()
    daily_df = monkey_info["daily_df"].copy()
    
    # SLIDE 1: Raw data with sequential day numbers
    print("  Creating Slide 1: Raw Data - Sequential Days")
    fig1 = create_raw_plot(monkey_name, species_name, session_df, daily_df, use_dates=False)
    plt.show()
    
    buf1 = io.BytesIO()
    fig1.savefig(buf1, format='png', dpi=150, bbox_inches='tight')
    buf1.seek(0)
    
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.shapes.add_picture(buf1, Inches(0.5), Inches(0.5), width=Inches(12.3))
    
    buf1.close()
    plt.close(fig1)
    
    # SLIDE 2: Smoothed data with sequential day numbers
    print("  Creating Slide 2: Smoothed Data - Sequential Days")
    fig2 = create_smoothed_plot(monkey_name, species_name, session_df, daily_df, use_dates=False)
    plt.show()
    
    buf2 = io.BytesIO()
    fig2.savefig(buf2, format='png', dpi=150, bbox_inches='tight')
    buf2.seek(0)
    
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.shapes.add_picture(buf2, Inches(0.5), Inches(0.5), width=Inches(12.3))
    
    buf2.close()
    plt.close(fig2)
    
    # SLIDE 3: Raw data with dates (WITH GAPS - NO CONNECTING)
    print("  Creating Slide 3: Raw Data - Calendar Dates (with gaps, no connecting)")
    fig3 = create_raw_plot(monkey_name, species_name, session_df, daily_df, use_dates=True)
    plt.show()
    
    buf3 = io.BytesIO()
    fig3.savefig(buf3, format='png', dpi=150, bbox_inches='tight')
    buf3.seek(0)
    
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    slide3.shapes.add_picture(buf3, Inches(0.5), Inches(0.5), width=Inches(12.3))
    
    buf3.close()
    plt.close(fig3)
    
    # SLIDE 4: Smoothed data with dates (WITH GAPS - NO CONNECTING)
    print("  Creating Slide 4: Smoothed Data - Calendar Dates (with gaps, no connecting)")
    fig4 = create_smoothed_plot(monkey_name, species_name, session_df, daily_df, use_dates=True)
    plt.show()
    
    buf4 = io.BytesIO()
    fig4.savefig(buf4, format='png', dpi=150, bbox_inches='tight')
    buf4.seek(0)
    
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    slide4.shapes.add_picture(buf4, Inches(0.5), Inches(0.5), width=Inches(12.3))
    
    buf4.close()
    plt.close(fig4)

# Save PowerPoint
pptx_path = downloads_dir / "rolling_averages_comparison.pptx"
prs.save(pptx_path)

print(f"\n{'='*60}")
print(f"COMPLETE!")
print(f"{'='*60}")
print(f"PowerPoint saved to: {pptx_path}")
print(f"Total slides created: {len(all_monkeys) * 4}")
print(f"\nSlides per monkey (in order):")
print(f"  1. Raw Data - Sequential Days (continuous)")
print(f"  2. Smoothed Data - Sequential Days (continuous)")
print(f"  3. Raw Data - Calendar Dates (vertical labels, vertical grid, gaps)")
print(f"  4. Smoothed Data - Calendar Dates (vertical labels, vertical grid, gaps)")
print(f"{'='*60}")