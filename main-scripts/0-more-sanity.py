import os
import pickle
from pathlib import Path
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import io

# ============================================================================
# CONFIGURATION
# ============================================================================
directory = "/Users/similovesyou/Desktop/qts/simian-behavior/data/py"
pickle_rick = os.path.join(directory, "data.pickle")

plt.rcParams["svg.fonttype"] = "none"

# ============================================================================
# SPECIFY MONKEYS AND DAY RANGES HERE
# ============================================================================
MONKEY_DAY_RANGES = {
    # Format: "monkey_name": (start_day, end_day)
    "wotan": (55, 100),
    "olaf": (80, 150),
    "havanna": (80, 90),
    #"iron": (26, 30),
    "lassa": (600, 640),
    "patsy": (0, 80),
    #"horus": (75, 95),
    "faramir": (40, 90),
    # Add more monkeys as needed
}

# ============================================================================
# LOAD DATA
# ============================================================================
with open(pickle_rick, "rb") as handle:
    data = pickle.load(handle)

# Remove monkeys with RFID issues
monkeys_to_remove = {"joy", "jipsy"}
for species in data:
    for name in list(data[species].keys()):
        if name in monkeys_to_remove:
            del data[species][name]

# ============================================================================
# HIERARCHY DATA
# ============================================================================
tonkean_hierarchy_file = os.path.join(directory, "hierarchy/tonkean/elo_matrix.xlsx")
rhesus_hierarchy_file = os.path.join(directory, "hierarchy/rhesus/elo_matrix.xlsx")

tonkean_abrv_2_fn = {
    "abr": "abricot", "ala": "alaryc", "alv": "alvin", "anu": "anubis",
    "bar": "barnabe", "ber": "berenice", "ces": "cesar", "dor": "dory",
    "eri": "eric", "jea": "jeanne", "lad": "lady", "las": "lassa",
    "nem": "nema", "nen": "nenno", "ner": "nereis", "ola": "olaf",
    "olg": "olga", "oli": "olli", "pac": "patchouli", "pat": "patsy",
    "wal": "wallace", "wat": "walt", "wot": "wotan", "yan": "yang",
    "yak": "yannick", "yin": "yin", "yoh": "yoh", "ult": "ulysse",
    "fic": "ficelle", "gan": "gandhi", "hav": "havanna", "con": "controle",
    "gai": "gaia", "her": "hercules", "han": "hanouk", "hor": "horus",
    "imo": "imoen", "ind": "indigo", "iro": "iron", "isi": "isis",
    "joy": "joy", "jip": "jipsy",
}

rhesus_abrv_2_fn = {
    "the": "theoden", "any": "anyanka", "djo": "djocko", "vla": "vladimir",
    "yel": "yelena", "bor": "boromir", "far": "faramir", "yva": "yvan",
    "baa": "baal", "spl": "spliff", "ol": "ol", "nat": "natasha",
    "arw": "arwen", "eow": "eowyn", "mar": "mar"
}

tonkean_hierarchy = pd.read_excel(tonkean_hierarchy_file)
rhesus_hierarchy = pd.read_excel(rhesus_hierarchy_file)

tonkean_hierarchy = tonkean_hierarchy.rename(columns=tonkean_abrv_2_fn)
rhesus_hierarchy = rhesus_hierarchy.rename(columns=rhesus_abrv_2_fn)

hierarchies = {
    "tonkean": tonkean_hierarchy,
    "rhesus": rhesus_hierarchy
}

# ============================================================================
# BUILD ELO DICTIONARY
# ============================================================================
elo = {}
for species, hierarchy in hierarchies.items():
    hierarchy["Date"] = pd.to_datetime(hierarchy["Date"])
    for name in hierarchy.columns[1:]:
        if name in hierarchy:
            monkey_elo = hierarchy[["Date", name]].dropna()
            if not monkey_elo.empty:
                first_date = monkey_elo["Date"].min()
                cutoff_date = first_date + pd.DateOffset(months=2)
                filtered_elo = monkey_elo[monkey_elo["Date"] > cutoff_date]
                elo[name] = filtered_elo

# ============================================================================
# COMPUTE DAILY PROPORTIONS
# ============================================================================
def compute_daily_proportions(species, name, monkey_data, elo_dict):
    """Compute daily performance proportions - only days with data."""
    attempts = monkey_data["attempts"].copy()
    
    attempts["instant_begin"] = pd.to_datetime(attempts["instant_begin"], unit="ms", errors="coerce")
    attempts["instant_end"] = pd.to_datetime(attempts["instant_end"], unit="ms", errors="coerce")
    attempts["date"] = attempts["instant_begin"].dt.date
    
    # Filter by ELO period if available
    if name in elo_dict and not elo_dict[name].empty:
        first_elo = elo_dict[name]["Date"].min().date()
        last_elo = elo_dict[name]["Date"].max().date()
        
        attempts = attempts[
            (attempts["date"] >= first_elo) & 
            (attempts["date"] <= last_elo)
        ]
    
    if attempts.empty:
        return pd.DataFrame()
    
    # Group by date
    daily_data = []
    day_counter = 1
    
    for date, day_attempts in attempts.groupby("date"):
        total_attempts = len(day_attempts)
        outcome_counts = day_attempts["result"].value_counts().to_dict()
        
        daily_data.append({
            "date": pd.to_datetime(date),
            "day_number": day_counter,
            "p_success": outcome_counts.get("success", 0.0) / total_attempts,
            "p_error": outcome_counts.get("error", 0.0) / total_attempts,
            "p_omission": outcome_counts.get("stepomission", 0.0) / total_attempts,
            "p_premature": outcome_counts.get("prematured", 0.0) / total_attempts,
            "total_attempts": total_attempts,
        })
        
        day_counter += 1
    
    daily_df = pd.DataFrame(daily_data)
    daily_df = daily_df.sort_values("date").reset_index(drop=True)
    
    return daily_df

# ============================================================================
# PLOTTING FUNCTION
# ============================================================================
def plot_day_range(monkey_name, species_name, daily_df, start_day, end_day):
    """Plot specific day range for a monkey."""
    
    # Filter to specified day range
    filtered_df = daily_df[
        (daily_df["day_number"] >= start_day) & 
        (daily_df["day_number"] <= end_day)
    ].copy()
    
    if filtered_df.empty:
        print(f"  WARNING: No data for days {start_day}-{end_day}")
        return None
    
    # Get date range
    date_start = filtered_df["date"].min().strftime("%Y-%m-%d")
    date_end = filtered_df["date"].max().strftime("%Y-%m-%d")
    actual_start_day = filtered_df["day_number"].min()
    actual_end_day = filtered_df["day_number"].max()
    
    print(f"\n  Day range: {actual_start_day} to {actual_end_day}")
    print(f"  Date range: {date_start} to {date_end}")
    
    # Prepare stacked data
    stacked_df = filtered_df[["p_success", "p_error", "p_premature", "p_omission"]].copy()
    stacked_df = stacked_df.div(stacked_df.sum(axis=1), axis=0)
    
    # Create plot
    fig, ax = plt.subplots(figsize=(12, 6))
    
    day_numbers = filtered_df["day_number"].values.astype(int)
    
    ax.stackplot(
        day_numbers,
        stacked_df["p_success"],
        stacked_df["p_error"],
        stacked_df["p_premature"],
        stacked_df["p_omission"],
        colors=["#5386b6", "#d57459", "#e8bc60", "#8d5993"],
        edgecolor="none", 
        linewidth=0, 
        baseline="zero",
        labels=["Success", "Error", "Premature", "Omission"],
        antialiased=False
    )
    
    # Add mean line
    overall_p_success = filtered_df["p_success"].mean()
    ax.axhline(y=overall_p_success, color="black", linestyle=(0, (8, 4)), linewidth=1.0)
    
    # Formatting
    ax.set_title(
        f"{monkey_name.capitalize()} ({species_name.capitalize()}) - Days {actual_start_day}-{actual_end_day}\n"
        f"Dates: {date_start} to {date_end}",
        fontsize=16, 
        fontweight="bold", 
        fontname="Times New Roman"
    )
    ax.set_xlabel("Day Number (consecutive)", fontsize=14, fontname="Times New Roman")
    ax.set_ylabel("Stacked Proportion", fontsize=14, fontname="Times New Roman")
    ax.set_xlim(day_numbers.min(), day_numbers.max())
    ax.set_ylim(0, 1)
    ax.set_yticks([0, 1])
    ax.legend(loc="upper left", fontsize=10, prop={"family": "Times New Roman"})
    ax.grid(True, linestyle="--", alpha=0.5)
    ax.xaxis.set_major_locator(plt.MaxNLocator(nbins=10, integer=True))
    
    plt.tight_layout()
    
    return fig, date_start, date_end

# ============================================================================
# MAIN PROCESSING
# ============================================================================
print("\nProcessing specified monkeys...")
print("="*60)

# Create PowerPoint
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Store summary data
summary_data = []

for monkey_name, (start_day, end_day) in MONKEY_DAY_RANGES.items():
    print(f"\nMonkey: {monkey_name}")
    print(f"Requested day range: {start_day} to {end_day}")
    
    # Find monkey in data
    found = False
    for species in ["tonkean", "rhesus"]:
        if species in data and monkey_name in data[species]:
            found = True
            species_name = species
            monkey_data = data[species][monkey_name]
            break
    
    if not found:
        print(f"  ERROR: Monkey '{monkey_name}' not found in data")
        continue
    
    # Compute daily proportions
    daily_df = compute_daily_proportions(species_name, monkey_name, monkey_data, elo)
    
    if daily_df.empty:
        print(f"  ERROR: No valid data for {monkey_name}")
        continue
    
    total_days = len(daily_df)
    print(f"  Total days available: {total_days}")
    
    if end_day > total_days:
        print(f"  WARNING: Requested end day {end_day} exceeds available days {total_days}")
        print(f"  Adjusting to maximum available: {total_days}")
        end_day = total_days
    
    # Create plot
    result = plot_day_range(monkey_name, species_name, daily_df, start_day, end_day)
    
    if result:
        fig, date_start, date_end = result
        plt.show()
        
        # Save to PowerPoint
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
        buf.seek(0)
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(buf, Inches(0.5), Inches(0.5), width=Inches(12.3))
        
        buf.close()
        plt.close(fig)
        
        # Store summary info
        summary_data.append({
            "Monkey": monkey_name.capitalize(),
            "Date Start": date_start,
            "Date End": date_end
        })
        
        print(f"  Added to PowerPoint")

# Save PowerPoint
output_dir = Path.home() / "Downloads"
pptx_path = output_dir / "monkey_day_ranges.pptx"
prs.save(pptx_path)
print(f"\nPowerPoint saved to: {pptx_path}")

# Create and save summary table
summary_df = pd.DataFrame(summary_data)
table_path = output_dir / "monkey_day_ranges_summary.csv"
summary_df.to_csv(table_path, index=False)
print(f"Summary table saved to: {table_path}")

print("\n" + "="*60)
print("SUMMARY TABLE:")
print("="*60)
print(summary_df.to_string(index=False))

print("\n" + "="*60)
print("COMPLETE!")
print(f"Total slides created: {len(summary_data)}")
print("="*60)