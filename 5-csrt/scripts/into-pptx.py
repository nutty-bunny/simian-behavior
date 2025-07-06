import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image  # Pillow is required

# Define paths
root_dir = "/Users/similovesyou/Desktop/qts/simian-behavior/visualisations"
output_pptx = "/Users/similovesyou/Downloads/all-up-to-date-plots.pptx"

# Slide size (in inches)
slide_width_in = 10  # usable width
slide_height_in = 7.5  # usable height
dpi = 96  # default dots per inch if unknown

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[6]

# Collect PNGs recursively
for dirpath, _, filenames in os.walk(root_dir):
    for file in sorted(filenames):  # sort alphabetically
        if file.lower().endswith(".png"):
            image_path = os.path.join(dirpath, file)

            # Open image to get size in inches
            with Image.open(image_path) as img:
                width_px, height_px = img.size
                width_in = width_px / dpi
                height_in = height_px / dpi

            # Calculate scaling factor
            scale_w = slide_width_in / width_in
            scale_h = slide_height_in / height_in
            scale = min(scale_w, scale_h)

            # Final size in inches
            final_width = width_in * scale
            final_height = height_in * scale

            # Centering
            left = Inches((slide_width_in - final_width) / 2)
            top = Inches((slide_height_in - final_height) / 2)

            # Add slide and image
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(image_path, left, top, width=Inches(final_width), height=Inches(final_height))

# Save presentation
prs.save(output_pptx)
print(f"PPTX saved to: {output_pptx}")
